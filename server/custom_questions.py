"""
IM Creator - Custom Questions Manager
Version: 8.0.0

Implements Requirement #10: Custom Questions → Extra Slides
Manages custom questions and renders them with professional templates
"""

from typing import Dict, List, Optional
from pptx.util import Inches, Pt
from models import CUSTOM_QUESTION_TEMPLATES


class CustomQuestionManager:
    """
    Manages custom questions and their slide rendering
    
    Features:
    - 6 professional slide templates
    - Template-specific rendering
    - User-friendly configuration
    """
    
    def __init__(self):
        self.templates = CUSTOM_QUESTION_TEMPLATES
    
    def create_custom_question(self, phase_id: str, config: Dict) -> Dict:
        """
        Create a custom question configuration
        
        Args:
            phase_id: Which questionnaire phase to add to
            config: {
                "label": "Question text",
                "type": "text|textarea|select",
                "slideTemplate": "text-only|two-column|...",
                "slideConfig": {...}  # Template-specific config
            }
        
        Returns:
            Complete question object ready for questionnaire
        """
        import time
        
        question_id = f"custom_{int(time.time() * 1000)}"
        
        question = {
            "id": question_id,
            "phase": phase_id,
            "label": config["label"],
            "type": config.get("type", "textarea"),
            "is_custom": True,
            "slide_template": config.get("slideTemplate", "text-only"),
            "slide_config": config.get("slideConfig", {}),
            "order": 999  # Place at end of phase
        }
        
        return question
    
    def render_custom_slide(self, slide, colors, data, question_config, layout_rec, context):
        """
        Render custom slide based on selected template
        
        Args:
            slide: python-pptx slide object
            colors: Theme colors
            data: Form data
            question_config: Custom question configuration
            layout_rec: AI layout recommendations
            context: Rendering context
        """
        template = question_config.get("slide_template", "text-only")
        field_id = question_config["id"]
        content = data.get(field_id, "")
        
        # Route to appropriate template renderer
        if template == "text-only":
            return self._render_text_only(slide, colors, question_config, content, layout_rec)
        
        elif template == "two-column":
            return self._render_two_column(slide, colors, question_config, content, layout_rec)
        
        elif template == "metrics-grid":
            return self._render_metrics_grid(slide, colors, question_config, content, layout_rec)
        
        elif template == "image-text":
            return self._render_image_text(slide, colors, question_config, content, layout_rec)
        
        elif template == "comparison":
            return self._render_comparison(slide, colors, question_config, content, layout_rec)
        
        elif template == "timeline":
            return self._render_timeline(slide, colors, question_config, content, layout_rec)
        
        else:
            # Fallback to text-only
            return self._render_text_only(slide, colors, question_config, content, layout_rec)
    
    def _render_text_only(self, slide, colors, config, content, layout_rec):
        """Render full-width text slide"""
        from pptx_generator import add_slide_header, smart_add_text, add_section_box
        
        title = config["label"]
        add_slide_header(slide, colors, title)
        
        add_section_box(slide, colors, 0.2, 0.9, 9.6, 4.0)
        smart_add_text(slide, 0.4, 1.2, 9.2, 3.4, content, 12, colors)
    
    def _render_two_column(self, slide, colors, config, content, layout_rec):
        """Render two-column split layout"""
        from pptx_generator import add_slide_header, smart_add_text, add_section_box
        
        title = config["label"]
        add_slide_header(slide, colors, title)
        
        # Parse content: "Left content|Right content"
        parts = content.split("|") if "|" in content else [content, ""]
        left_content = parts[0].strip() if len(parts) > 0 else ""
        right_content = parts[1].strip() if len(parts) > 1 else ""
        
        # Get titles from config
        left_title = config.get("slide_config", {}).get("left_title", "Left")
        right_title = config.get("slide_config", {}).get("right_title", "Right")
        
        # Left column
        add_section_box(slide, colors, 0.2, 0.9, 4.7, 4.0, left_title)
        smart_add_text(slide, 0.4, 1.5, 4.3, 3.0, left_content, 11, colors)
        
        # Right column
        add_section_box(slide, colors, 5.1, 0.9, 4.7, 4.0, right_title, colors["secondary"])
        smart_add_text(slide, 5.3, 1.5, 4.3, 3.0, right_content, 11, colors)
    
    def _render_metrics_grid(self, slide, colors, config, content, layout_rec):
        """Render 2x2 or 3x3 grid of metrics"""
        from pptx_generator import add_slide_header, add_metric_card
        
        title = config["label"]
        add_slide_header(slide, colors, title)
        
        # Parse content: "Revenue: $50M|Growth: 25%|Margin: 40%|ARR: $30M"
        metric_items = content.split("|")
        metrics = []
        for item in metric_items:
            if ":" in item:
                label, value = item.split(":", 1)
                metrics.append({"label": label.strip(), "value": value.strip()})
        
        # Positions for 2x2 grid
        positions = [
            (0.2, 1.2, 4.7, 1.3),   # Top-left
            (5.1, 1.2, 4.7, 1.3),   # Top-right
            (0.2, 2.7, 4.7, 1.3),   # Bottom-left
            (5.1, 2.7, 4.7, 1.3)    # Bottom-right
        ]
        
        # Render up to 4 metrics
        for i, metric in enumerate(metrics[:4]):
            x, y, w, h = positions[i]
            add_metric_card(slide, colors, x, y, w, h, metric["value"], metric["label"])
    
    # Additional template renderers would go here...
    # _render_image_text, _render_comparison, _render_timeline