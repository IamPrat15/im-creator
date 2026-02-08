"""
IM Creator Python Backend - PPTX Generator
Version: 8.3.0

Implements Requirements #15-18 + Deloitte-Quality Design Enhancements:
- Universal createSlide() wrapper
- Dedicated render functions for each slide type
- Chart helper addChartByType()
- AI-powered layout recommendations applied consistently

v8.3.0 Enhancements (Deloitte-Quality Design):
- Section divider slides (dark full-bleed backgrounds)
- Enhanced case study layout (sidebar card + two-column projects + metric results)
- Large infographic-style metric cards (big number + label + icon indicator)
- Proper stacked bar charts with CAGR annotations
- Enhanced timeline with horizontal arrow + alternating milestone labels
- Leadership org-chart style layout with role tiers
- Table of Contents with section numbering
- Company Overview with rich metric row
- Fixed risk factors renderer (blank layout compatible)
- Missing slide type handlers added (toc, company-overview, leadership, risks)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from typing import Dict, List, Optional
from models import DESIGN, INDUSTRY_DATA, DOCUMENT_CONFIGS, get_theme_colors
from utils import (
    truncate_text, truncate_description, format_currency, format_date,
    parse_lines, parse_pipe_separated, calculate_cagr, safe_float, safe_int,
    adjusted_font, extract_percentage,
    get_slides_for_document_type, get_buyer_specific_content, get_industry_specific_content
)
from ai_layout_engine import analyze_data_for_layout_sync

# ============================================================================
# COLOR HELPER
# ============================================================================
def hex_to_rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

# ============================================================================
# CHART DISPATCHER
# ============================================================================
def add_chart_by_type(slide, colors, x, y, w, h, chart_type, data, font_adj=0):
    if not data or chart_type == "none":
        return None
    dispatch = {
        "bar": add_bar_chart, "pie": add_pie_chart, "donut": add_donut_chart,
        "timeline": add_timeline, "progress": add_progress_bars, "stacked-bar": add_stacked_bar_chart
    }
    fn = dispatch.get(chart_type)
    return fn(slide, colors, x, y, w, h, data, font_adj) if fn else None

# ============================================================================
# BASE SLIDE COMPONENTS
# ============================================================================
def add_slide_header(slide, colors, title, subtitle=None, font_adj=0):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(5.625))
    bg.fill.solid(); bg.fill.fore_color.rgb = hex_to_rgb(colors["white"]); bg.line.fill.background()
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.1), Inches(0.85))
    bar.fill.solid(); bar.fill.fore_color.rgb = hex_to_rgb(colors["secondary"]); bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.5))
    tb.text_frame.paragraphs[0].text = truncate_text(title, 80)
    tb.text_frame.paragraphs[0].font.size = Pt(adjusted_font(DESIGN["fonts"]["title"], font_adj))
    tb.text_frame.paragraphs[0].font.bold = True
    tb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["primary"])
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.3), Inches(0.62), Inches(9.4), Inches(0.22))
        sb.text_frame.paragraphs[0].text = subtitle
        sb.text_frame.paragraphs[0].font.size = Pt(adjusted_font(DESIGN["fonts"]["subtitle"], font_adj))
        sb.text_frame.paragraphs[0].font.italic = True
        sb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["text_light"])
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(0.88), Inches(9.4), Inches(0.02))
    line.fill.solid(); line.fill.fore_color.rgb = hex_to_rgb(colors["accent"]); line.line.fill.background()

def add_slide_footer(slide, colors, page_number):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.2), Inches(10), Inches(0.015))
    line.fill.solid(); line.fill.fore_color.rgb = hex_to_rgb(colors["primary"]); line.line.fill.background()
    cb = slide.shapes.add_textbox(Inches(0.3), Inches(5.28), Inches(3), Inches(0.22))
    cb.text_frame.paragraphs[0].text = "Strictly Private & Confidential"
    cb.text_frame.paragraphs[0].font.size = Pt(9); cb.text_frame.paragraphs[0].font.italic = True
    cb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["text_light"])
    nb = slide.shapes.add_textbox(Inches(9.2), Inches(5.28), Inches(0.5), Inches(0.22))
    nb.text_frame.paragraphs[0].text = str(page_number)
    nb.text_frame.paragraphs[0].font.size = Pt(10); nb.text_frame.paragraphs[0].font.bold = True
    nb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["primary"])
    nb.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

def add_section_box(slide, colors, x, y, w, h, title=None, title_bg=None, font_adj=0):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = hex_to_rgb(colors["light_bg"])
    box.line.color.rgb = hex_to_rgb(colors["border"])
    if title:
        hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.36))
        hdr.fill.solid(); hdr.fill.fore_color.rgb = hex_to_rgb(title_bg or colors["primary"]); hdr.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(x+0.12), Inches(y+0.02), Inches(w-0.24), Inches(0.32))
        tb.text_frame.paragraphs[0].text = truncate_text(title, 45)
        tb.text_frame.paragraphs[0].font.size = Pt(adjusted_font(DESIGN["fonts"]["section_header"], font_adj))
        tb.text_frame.paragraphs[0].font.bold = True
        tb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["white"])

def add_metric_card(slide, colors, x, y, w, h, value, label, font_adj=0):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid(); card.fill.fore_color.rgb = hex_to_rgb(colors["light_bg"])
    card.line.color.rgb = hex_to_rgb(colors["border"])
    vb = slide.shapes.add_textbox(Inches(x+0.08), Inches(y+0.08), Inches(w-0.16), Inches(h*0.55))
    vb.text_frame.paragraphs[0].text = str(value)
    vb.text_frame.paragraphs[0].font.size = Pt(adjusted_font(DESIGN["fonts"]["metric_medium"], font_adj))
    vb.text_frame.paragraphs[0].font.bold = True
    vb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["primary"])
    lb = slide.shapes.add_textbox(Inches(x+0.08), Inches(y+h*0.58), Inches(w-0.16), Inches(h*0.38))
    lb.text_frame.paragraphs[0].text = label
    lb.text_frame.paragraphs[0].font.size = Pt(adjusted_font(DESIGN["fonts"]["metric_label"], font_adj))
    lb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["text_light"])

# ============================================================================
# v8.3.0: LARGE INFOGRAPHIC METRIC CARD (Deloitte-style)
# ============================================================================
def add_metric_card_large(slide, colors, x, y, w, h, value, label, icon_char="●", font_adj=0):
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid(); card.fill.fore_color.rgb = hex_to_rgb(colors["light_bg"])
    card.line.color.rgb = hex_to_rgb(colors["border"])
    # Left accent
    acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.06), Inches(h))
    acc.fill.solid(); acc.fill.fore_color.rgb = hex_to_rgb(colors["secondary"]); acc.line.fill.background()
    # Icon circle
    isz = 0.3
    ic = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+0.15), Inches(y+(h-isz)/2), Inches(isz), Inches(isz))
    ic.fill.solid(); ic.fill.fore_color.rgb = hex_to_rgb(colors["secondary"]); ic.line.fill.background()
    itb = slide.shapes.add_textbox(Inches(x+0.15), Inches(y+(h-isz)/2), Inches(isz), Inches(isz))
    itb.text_frame.paragraphs[0].text = icon_char
    itb.text_frame.paragraphs[0].font.size = Pt(12)
    itb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["white"])
    itb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Value
    tx = x + 0.55; tw = w - 0.65
    vb = slide.shapes.add_textbox(Inches(tx), Inches(y+0.05), Inches(tw), Inches(h*0.55))
    vb.text_frame.paragraphs[0].text = str(value)
    vb.text_frame.paragraphs[0].font.size = Pt(adjusted_font(DESIGN["fonts"]["metric"], font_adj))
    vb.text_frame.paragraphs[0].font.bold = True
    vb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["primary"])
    # Label
    lb = slide.shapes.add_textbox(Inches(tx), Inches(y+h*0.55), Inches(tw), Inches(h*0.4))
    lb.text_frame.paragraphs[0].text = label
    lb.text_frame.paragraphs[0].font.size = Pt(adjusted_font(DESIGN["fonts"]["metric_label"], font_adj))
    lb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["text_light"])

def add_metric_row(slide, colors, metrics, y, x_start=0.3, total_width=9.4, font_adj=0):
    if not metrics: return
    num = min(len(metrics), 6); gap = 0.15
    cw = (total_width - (gap * (num-1))) / num; ch = 0.85
    icons = ["◆","▲","●","★","■","▶"]
    for i in range(num):
        v, l = metrics[i][0], metrics[i][1]
        ic = metrics[i][2] if len(metrics[i]) > 2 else icons[i % len(icons)]
        add_metric_card_large(slide, colors, x_start+(i*(cw+gap)), y, cw, ch, v, l, ic, font_adj)

# ============================================================================
# CHART FUNCTIONS
# ============================================================================
def add_bar_chart(slide, colors, x, y, w, h, data, font_adj=0):
    if not data: return
    cd = CategoryChartData()
    cd.categories = [d.get("label","") for d in data]
    cd.add_series("Values", tuple(safe_float(d.get("value",0)) for d in data))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), cd).chart
    chart.has_legend = False; chart.plots[0].has_data_labels = True
    for s in chart.series: s.format.fill.solid(); s.format.fill.fore_color.rgb = hex_to_rgb(colors["primary"])
    return chart

def add_pie_chart(slide, colors, x, y, w, h, data, font_adj=0):
    if not data: return
    cd = CategoryChartData()
    cd.categories = [d.get("label","") for d in data]
    cd.add_series("Values", tuple(safe_float(d.get("value",0)) for d in data))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(x), Inches(y), Inches(w), Inches(h), cd).chart
    chart.has_legend = True; chart.plots[0].has_data_labels = True
    return chart

def add_donut_chart(slide, colors, x, y, w, h, data, font_adj=0):
    if not data: return
    cd = CategoryChartData()
    cd.categories = [d.get("label","") for d in data]
    cd.add_series("Values", tuple(safe_float(d.get("value",0)) for d in data))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(x), Inches(y), Inches(w), Inches(h), cd).chart
    chart.has_legend = True; chart.plots[0].has_data_labels = True
    return chart

# v8.3.0: ENHANCED TIMELINE
def add_timeline(slide, colors, x, y, w, h, milestones, font_adj=0):
    if not milestones: return
    num = len(milestones)
    axis_y = y + h*0.45; axis_h = 0.06
    ab = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(axis_y), Inches(w), Inches(axis_h))
    ab.fill.solid(); ab.fill.fore_color.rgb = hex_to_rgb(colors["primary"]); ab.line.fill.background()
    # Arrow head
    ah = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x+w-0.05), Inches(axis_y-0.07), Inches(0.2), Inches(axis_h+0.14))
    ah.fill.solid(); ah.fill.fore_color.rgb = hex_to_rgb(colors["primary"]); ah.line.fill.background()
    step = w / max(num, 1); nsz = 0.28
    for i, ms in enumerate(milestones):
        mx = x + (i*step) + (step/2) - (nsz/2)
        # Node
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(mx), Inches(axis_y-(nsz-axis_h)/2), Inches(nsz), Inches(nsz))
        c.fill.solid(); c.fill.fore_color.rgb = hex_to_rgb(colors["secondary"]); c.line.fill.background()
        yt = str(ms.get("year", ms.get("label","")))[:6]
        ntb = slide.shapes.add_textbox(Inches(mx), Inches(axis_y-(nsz-axis_h)/2), Inches(nsz), Inches(nsz))
        ntb.text_frame.paragraphs[0].text = yt
        ntb.text_frame.paragraphs[0].font.size = Pt(adjusted_font(8, font_adj))
        ntb.text_frame.paragraphs[0].font.bold = True
        ntb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["white"])
        ntb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        # Alternating label
        lbl = truncate_text(ms.get("label", ms.get("event","")), 25)
        lw = step - 0.1; lx = x + (i*step) + 0.05
        ly = axis_y - 0.55 if i%2==0 else axis_y + nsz + 0.05
        # Connector
        cx = mx + nsz/2 - 0.01
        cy = (axis_y - 0.25) if i%2==0 else (axis_y + nsz)
        ch2 = 0.2 if i%2==0 else 0.15
        cn = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx), Inches(cy), Inches(0.02), Inches(ch2))
        cn.fill.solid(); cn.fill.fore_color.rgb = hex_to_rgb(colors["border"]); cn.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(lx), Inches(ly), Inches(lw), Inches(0.4))
        tb.text_frame.text = lbl
        tb.text_frame.paragraphs[0].font.size = Pt(adjusted_font(9, font_adj))
        tb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["text"])
        tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        tb.text_frame.word_wrap = True

def add_progress_bars(slide, colors, x, y, w, h, items, font_adj=0):
    if not items: return
    num = len(items); bh = min(0.25, (h-0.1)/max(num,1)); gap = 0.1
    for i, item in enumerate(items):
        by = y + (i*(bh+gap)); lbl = item.get("label",""); val = safe_float(item.get("value",0))
        bgb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(by), Inches(w), Inches(bh))
        bgb.fill.solid(); bgb.fill.fore_color.rgb = hex_to_rgb(colors["light_bg"]); bgb.line.fill.background()
        pw = max((w*val)/100, 0.3)
        pb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(by), Inches(pw), Inches(bh))
        pb.fill.solid(); pb.fill.fore_color.rgb = hex_to_rgb(colors["primary"]); pb.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(x+0.1), Inches(by+0.05), Inches(w-0.2), Inches(bh-0.1))
        tb.text_frame.text = f"{lbl}: {val}%"
        tb.text_frame.paragraphs[0].font.size = Pt(adjusted_font(11, font_adj))
        tb.text_frame.paragraphs[0].font.bold = True
        tb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["white"])

# v8.3.0: PROPER STACKED BAR
def add_stacked_bar_chart(slide, colors, x, y, w, h, data, font_adj=0):
    if not data: return
    cd = CategoryChartData()
    cd.categories = [d.get("label","") for d in data]
    cd.add_series("Services", tuple(safe_float(d.get("value",0)) for d in data))
    has2 = any(d.get("value2") for d in data)
    if has2: cd.add_series("License/Resale", tuple(safe_float(d.get("value2",0)) for d in data))
    ct = XL_CHART_TYPE.COLUMN_STACKED if has2 else XL_CHART_TYPE.COLUMN_CLUSTERED
    chart = slide.shapes.add_chart(ct, Inches(x), Inches(y), Inches(w), Inches(h), cd).chart
    chart.has_legend = has2; chart.plots[0].has_data_labels = True
    cc = colors.get("chart_colors", [colors["primary"], colors["secondary"]])
    for i, s in enumerate(chart.series):
        s.format.fill.solid(); s.format.fill.fore_color.rgb = hex_to_rgb(cc[i % len(cc)])
    return chart

def add_cagr_annotation(slide, colors, x, y, w, cagr_value, font_adj=0):
    if not cagr_value: return
    al = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x+0.2), Inches(y), Inches(w-0.4), Inches(0.03))
    al.fill.solid(); al.fill.fore_color.rgb = hex_to_rgb(colors["secondary"]); al.line.fill.background()
    for ax in [x+0.15, x+w-0.25]:
        d = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(ax), Inches(y-0.04), Inches(0.1), Inches(0.1))
        d.fill.solid(); d.fill.fore_color.rgb = hex_to_rgb(colors["secondary"]); d.line.fill.background()
    ct = slide.shapes.add_textbox(Inches(x+w*0.25), Inches(y-0.22), Inches(w*0.5), Inches(0.2))
    ct.text_frame.paragraphs[0].text = f"CAGR: {cagr_value}%"
    ct.text_frame.paragraphs[0].font.size = Pt(adjusted_font(9, font_adj))
    ct.text_frame.paragraphs[0].font.bold = True
    ct.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["secondary"])
    ct.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ============================================================================
# v8.3.0: SECTION DIVIDER (Deloitte dark full-bleed)
# ============================================================================
def render_section_divider(slide, colors, section_title, section_number=None, font_adj=0):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(5.625))
    bg.fill.solid(); bg.fill.fore_color.rgb = hex_to_rgb(colors["primary"]); bg.line.fill.background()
    al = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.2), Inches(8), Inches(0.04))
    al.fill.solid(); al.fill.fore_color.rgb = hex_to_rgb(colors["secondary"]); al.line.fill.background()
    if section_number is not None:
        bsz = 0.6
        b = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(1.3), Inches(bsz), Inches(bsz))
        b.fill.solid(); b.fill.fore_color.rgb = hex_to_rgb(colors["secondary"]); b.line.fill.background()
        btb = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(bsz), Inches(bsz))
        btb.text_frame.paragraphs[0].text = str(section_number)
        btb.text_frame.paragraphs[0].font.size = Pt(22); btb.text_frame.paragraphs[0].font.bold = True
        btb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["white"])
        btb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    ttb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.0))
    ttb.text_frame.paragraphs[0].text = section_title
    ttb.text_frame.paragraphs[0].font.size = Pt(40); ttb.text_frame.paragraphs[0].font.bold = True
    ttb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["white"])
    stb = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.4))
    stb.text_frame.paragraphs[0].text = "Strictly Private & Confidential"
    stb.text_frame.paragraphs[0].font.size = Pt(12); stb.text_frame.paragraphs[0].font.italic = True
    stb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["white"])


# ============================================================================
# RENDER FUNCTIONS
# ============================================================================
def render_executive_summary(slide, colors, data, page_num, layout_rec, context):
    font_adj = layout_rec.get("font_adjustment", 0)
    chart_type = layout_rec.get("chart_type", "bar")
    layout = layout_rec.get("layout", "two-column")
    buyer_types = data.get("targetBuyerType") or data.get("target_buyer_type") or ["strategic"]
    vertical = data.get("primaryVertical") or data.get("primary_vertical") or "technology"
    buyer_content = get_buyer_specific_content(buyer_types, "executive-summary", data)
    industry_content = get_industry_specific_content(vertical, "executive-summary")
    add_slide_header(slide, colors, "Executive Summary", industry_content.get("context"), font_adj)
    add_slide_footer(slide, colors, page_num)
    if layout == "two-column":
        add_section_box(slide, colors, 0.3, 0.95, 4.5, 2.8, "Company Overview", font_adj=font_adj)
        desc = data.get("companyDescription") or data.get("company_description") or ""
        tb = slide.shapes.add_textbox(Inches(0.45), Inches(1.45), Inches(4.2), Inches(2.0))
        tb.text_frame.text = truncate_description(desc, 300)
        tb.text_frame.paragraphs[0].font.size = Pt(adjusted_font(DESIGN["fonts"]["body"], font_adj))
        tb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["text"])
        tb.text_frame.word_wrap = True
        add_section_box(slide, colors, 5.0, 0.95, 4.7, 2.8, "Revenue Growth", colors["secondary"], font_adj)
        revenue_data = []
        for key, label in [("revenueFY24","FY24"),("revenueFY25","FY25"),("revenueFY26P","FY26P"),("revenueFY27P","FY27P")]:
            val = safe_float(data.get(key))
            if val: revenue_data.append({"label": label, "value": val})
        if revenue_data and chart_type != "none":
            if len(revenue_data) >= 2:
                cagr = calculate_cagr(revenue_data[0]["value"], revenue_data[-1]["value"], len(revenue_data)-1)
                if cagr: add_cagr_annotation(slide, colors, 5.2, 1.35, 4.3, round(cagr,1), font_adj)
            add_chart_by_type(slide, colors, 5.2, 1.5, 4.3, 2.0, chart_type, revenue_data, font_adj)
        # Bottom metric row
        metrics = []
        metrics.append((str(data.get("foundedYear") or "N/A"), "Founded", "◆"))
        metrics.append((str(data.get("employeeCountFT") or "N/A"), "Employees", "▲"))
        t10 = data.get("top10Concentration") or data.get("topClientCount") or ""
        if t10: metrics.append((f"{t10}%", "Top 10 Conc.", "●"))
        else: metrics.append((str(data.get("headquarters") or "N/A")[:15], "Headquarters", "●"))
        ebitda = data.get("ebitdaMarginFY25") or data.get("ebitda_margin_fy25") or ""
        if ebitda: metrics.append((f"{ebitda}%", "EBITDA Margin", "★"))
        add_metric_row(slide, colors, metrics, y=4.0, font_adj=font_adj)
    elif layout == "full-width":
        add_section_box(slide, colors, 0.3, 0.95, 9.4, 3.8, "Executive Summary", font_adj=font_adj)
        desc = data.get("companyDescription") or ""
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9.0), Inches(2.8))
        tb.text_frame.text = truncate_description(desc, 600)
        tb.text_frame.paragraphs[0].font.size = Pt(adjusted_font(DESIGN["fonts"]["body_large"], font_adj))
        tb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["text"])

def render_services(slide, colors, data, page_num, layout_rec, context):
    font_adj = layout_rec.get("font_adjustment", 0)
    chart_type = layout_rec.get("chart_type", "donut")
    add_slide_header(slide, colors, "Service Lines & Capabilities", font_adj=font_adj)
    add_slide_footer(slide, colors, page_num)
    service_text = data.get("serviceLines") or data.get("service_lines") or ""
    services = parse_pipe_separated(service_text, 8)
    add_section_box(slide, colors, 0.3, 0.95, 4.8, 3.8, "Service Offerings", font_adj=font_adj)
    ind_colors = [colors["primary"], colors["secondary"], colors["accent"], colors.get("success","38A169")]
    y_pos = 1.5
    for idx, svc in enumerate(services[:6]):
        if len(svc) >= 2:
            name = truncate_text(svc[0], 30); pct = svc[1] if len(svc)>1 else ""
            desc = svc[2] if len(svc)>2 else ""
            ic = ind_colors[idx % len(ind_colors)]
            ind = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y_pos+0.02), Inches(0.12), Inches(0.12))
            ind.fill.solid(); ind.fill.fore_color.rgb = hex_to_rgb(ic); ind.line.fill.background()
            tb = slide.shapes.add_textbox(Inches(0.72), Inches(y_pos-0.02), Inches(4.2), Inches(0.25))
            tb.text_frame.text = f"{name} ({pct})"
            tb.text_frame.paragraphs[0].font.size = Pt(adjusted_font(DESIGN["fonts"]["body"], font_adj))
            tb.text_frame.paragraphs[0].font.bold = True
            tb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["text"])
            if desc:
                dtb = slide.shapes.add_textbox(Inches(0.72), Inches(y_pos+0.2), Inches(4.2), Inches(0.25))
                dtb.text_frame.text = truncate_text(desc, 50)
                dtb.text_frame.paragraphs[0].font.size = Pt(adjusted_font(DESIGN["fonts"]["body_small"], font_adj))
                dtb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["text_light"])
                y_pos += 0.55
            else: y_pos += 0.4
    add_section_box(slide, colors, 5.3, 0.95, 4.4, 3.8, "Revenue by Service", colors["secondary"], font_adj)
    chart_data = []
    for svc in services:
        if len(svc) >= 2:
            n = truncate_text(svc[0],20); p = safe_float(svc[1].replace("%","").strip())
            if p > 0: chart_data.append({"label": n, "value": p})
    if chart_data and chart_type != "none":
        add_chart_by_type(slide, colors, 5.5, 1.5, 4.0, 2.8, chart_type, chart_data, font_adj)

def render_clients(slide, colors, data, page_num, layout_rec, context):
    font_adj = layout_rec.get("font_adjustment", 0)
    add_slide_header(slide, colors, "Client Portfolio", font_adj=font_adj)
    add_slide_footer(slide, colors, page_num)
    client_text = data.get("topClients") or data.get("top_clients") or ""
    clients = parse_pipe_separated(client_text, 12)
    add_section_box(slide, colors, 0.3, 0.95, 3.2, 3.8, "Key Metrics", font_adj=font_adj)
    top10 = data.get("top10Concentration") or data.get("top_10_concentration") or "N/A"
    nrr = data.get("netRetention") or data.get("net_retention") or "N/A"
    add_metric_card(slide, colors, 0.45, 1.5, 2.9, 0.75, f"{top10}%", "Top 10 Concentration", font_adj)
    add_metric_card(slide, colors, 0.45, 2.5, 2.9, 0.75, f"{nrr}%", "Net Revenue Retention", font_adj)
    add_section_box(slide, colors, 3.7, 0.95, 6.0, 3.8, "Top Clients", colors["secondary"], font_adj)
    y_pos = 1.5
    for cl in clients[:10]:
        if cl:
            n = truncate_text(cl[0] if len(cl)>0 else "", 40)
            tb = slide.shapes.add_textbox(Inches(3.9), Inches(y_pos), Inches(5.6), Inches(0.3))
            tb.text_frame.text = f"• {n}"
            tb.text_frame.paragraphs[0].font.size = Pt(adjusted_font(DESIGN["fonts"]["body_small"], font_adj))
            tb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["text"])
            y_pos += 0.32

def render_financials(slide, colors, data, page_num, layout_rec, context):
    font_adj = layout_rec.get("font_adjustment", 0)
    chart_type = layout_rec.get("chart_type", "bar")
    add_slide_header(slide, colors, "Financial Performance", font_adj=font_adj)
    add_slide_footer(slide, colors, page_num)
    add_section_box(slide, colors, 0.3, 0.95, 4.5, 3.8, "Revenue Trend (INR Cr)", font_adj=font_adj)
    revenue_data = []
    for key, label in [("revenueFY24","FY24"),("revenueFY25","FY25"),("revenueFY26P","FY26P"),("revenueFY27P","FY27P")]:
        val = safe_float(data.get(key))
        if val: revenue_data.append({"label": label, "value": val})
    if revenue_data and chart_type != "none":
        if len(revenue_data) >= 2:
            cagr = calculate_cagr(revenue_data[0]["value"], revenue_data[-1]["value"], len(revenue_data)-1)
            if cagr: add_cagr_annotation(slide, colors, 0.5, 1.35, 4.1, round(cagr,1), font_adj)
        add_chart_by_type(slide, colors, 0.5, 1.5, 4.1, 2.8, chart_type, revenue_data, font_adj)
    add_section_box(slide, colors, 5.0, 0.95, 4.7, 3.8, "Profitability Metrics", colors["secondary"], font_adj)
    ebitda = data.get("ebitdaMarginFY25") or data.get("ebitda_margin_fy25") or "N/A"
    gross = data.get("grossMargin") or data.get("gross_margin") or "N/A"
    net = data.get("netProfitMargin") or data.get("net_profit_margin") or "N/A"
    add_metric_card(slide, colors, 5.2, 1.5, 4.3, 0.65, f"{ebitda}%", "EBITDA Margin FY25", font_adj)
    add_metric_card(slide, colors, 5.2, 2.35, 4.3, 0.65, f"{gross}%", "Gross Margin", font_adj)
    add_metric_card(slide, colors, 5.2, 3.2, 4.3, 0.65, f"{net}%", "Net Profit Margin", font_adj)

# v8.3.0: ENHANCED CASE STUDY (Deloitte sidebar + sections)
def render_case_study(slide, colors, data, page_num, case_study, layout_rec, context):
    font_adj = layout_rec.get("font_adjustment", 0)
    client = case_study.get("client", "Client")
    add_slide_header(slide, colors, f"Case Study: {truncate_text(client, 50)}", font_adj=font_adj)
    add_slide_footer(slide, colors, page_num)
    # LEFT SIDEBAR
    sx, sy, sw, sh = 0.3, 0.95, 2.8, 4.0
    sb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(sx), Inches(sy), Inches(sw), Inches(sh))
    sb.fill.solid(); sb.fill.fore_color.rgb = hex_to_rgb(colors["light_bg"]); sb.line.color.rgb = hex_to_rgb(colors["border"])
    shdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(sx), Inches(sy), Inches(sw), Inches(0.36))
    shdr.fill.solid(); shdr.fill.fore_color.rgb = hex_to_rgb(colors["primary"]); shdr.line.fill.background()
    htb = slide.shapes.add_textbox(Inches(sx+0.12), Inches(sy+0.02), Inches(sw-0.24), Inches(0.32))
    htb.text_frame.paragraphs[0].text = truncate_text(client, 25)
    htb.text_frame.paragraphs[0].font.size = Pt(adjusted_font(12, font_adj))
    htb.text_frame.paragraphs[0].font.bold = True
    htb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["white"])
    # Metadata
    my = sy + 0.5
    meta = []
    for k, lbl in [("industry","Industry"),("customerSince","Customer Since"),("engagementType","Engagement"),("platform","Platform"),("pricingModel","Pricing Model")]:
        v = case_study.get(k, "")
        if v: meta.append((lbl, v))
    if not meta: meta = [("Type","Enterprise"),("Engagement","Multi-year")]
    for ml, mv in meta:
        mlb = slide.shapes.add_textbox(Inches(sx+0.15), Inches(my), Inches(sw-0.3), Inches(0.18))
        mlb.text_frame.paragraphs[0].text = ml
        mlb.text_frame.paragraphs[0].font.size = Pt(8); mlb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["text_light"])
        mvb = slide.shapes.add_textbox(Inches(sx+0.15), Inches(my+0.15), Inches(sw-0.3), Inches(0.2))
        mvb.text_frame.paragraphs[0].text = truncate_text(str(mv), 25)
        mvb.text_frame.paragraphs[0].font.size = Pt(10); mvb.text_frame.paragraphs[0].font.bold = True
        mvb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["text"])
        my += 0.4
    # RIGHT: Challenge/Solution/Results
    rx = sx + sw + 0.2; rw = 9.4 - sw - 0.2
    sections = [
        ("Challenge", case_study.get("challenge",""), colors["accent"]),
        ("Solution", case_study.get("solution",""), colors["primary"]),
        ("Results", case_study.get("results",""), colors["secondary"])
    ]
    sec_y = sy; sec_h = sh / 3 - 0.05
    for title, content, hc in sections:
        sh2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(rx), Inches(sec_y), Inches(rw), Inches(0.3))
        sh2.fill.solid(); sh2.fill.fore_color.rgb = hex_to_rgb(hc); sh2.line.fill.background()
        stb = slide.shapes.add_textbox(Inches(rx+0.1), Inches(sec_y+0.01), Inches(rw-0.2), Inches(0.28))
        stb.text_frame.paragraphs[0].text = title
        stb.text_frame.paragraphs[0].font.size = Pt(adjusted_font(11, font_adj))
        stb.text_frame.paragraphs[0].font.bold = True
        stb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["white"])
        cbg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(rx), Inches(sec_y+0.3), Inches(rw), Inches(sec_h-0.3))
        cbg.fill.solid(); cbg.fill.fore_color.rgb = hex_to_rgb(colors["light_bg"])
        cbg.line.color.rgb = hex_to_rgb(colors["border"])
        ctb = slide.shapes.add_textbox(Inches(rx+0.1), Inches(sec_y+0.35), Inches(rw-0.2), Inches(sec_h-0.4))
        ctb.text_frame.text = truncate_description(content, 180)
        ctb.text_frame.paragraphs[0].font.size = Pt(adjusted_font(DESIGN["fonts"]["body_small"], font_adj))
        ctb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["text"])
        ctb.text_frame.word_wrap = True
        sec_y += sec_h + 0.05

def render_growth(slide, colors, data, page_num, layout_rec, context):
    font_adj = layout_rec.get("font_adjustment", 0)
    add_slide_header(slide, colors, "Growth Strategy & Roadmap", font_adj=font_adj)
    add_slide_footer(slide, colors, page_num)
    add_section_box(slide, colors, 0.3, 0.95, 4.5, 3.8, "Key Growth Drivers", font_adj=font_adj)
    drivers = parse_lines(data.get("growthDrivers") or data.get("growth_drivers") or "", 6)
    y_pos = 1.5
    for d in drivers:
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(y_pos+0.02), Inches(0.22), Inches(0.22))
        c.fill.solid(); c.fill.fore_color.rgb = hex_to_rgb(colors["secondary"]); c.line.fill.background()
        itb = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos+0.02), Inches(0.22), Inches(0.22))
        itb.text_frame.paragraphs[0].text = "▶"
        itb.text_frame.paragraphs[0].font.size = Pt(8)
        itb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["white"])
        itb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        tb = slide.shapes.add_textbox(Inches(0.82), Inches(y_pos), Inches(3.8), Inches(0.3))
        tb.text_frame.text = truncate_text(d, 50)
        tb.text_frame.paragraphs[0].font.size = Pt(adjusted_font(DESIGN["fonts"]["body"], font_adj))
        tb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["text"])
        y_pos += 0.45
    add_section_box(slide, colors, 5.0, 0.95, 4.7, 3.8, "Strategic Goals", colors["secondary"], font_adj)
    short_goals = parse_lines(data.get("shortTermGoals") or data.get("short_term_goals") or "", 3)
    medium_goals = parse_lines(data.get("mediumTermGoals") or data.get("medium_term_goals") or "", 3)
    y_pos = 1.5
    if short_goals:
        tb = slide.shapes.add_textbox(Inches(5.2), Inches(y_pos), Inches(4.3), Inches(0.25))
        tb.text_frame.text = "Short-Term (0-12 months)"
        tb.text_frame.paragraphs[0].font.size = Pt(adjusted_font(DESIGN["fonts"]["body_large"], font_adj))
        tb.text_frame.paragraphs[0].font.bold = True
        tb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["primary"])
        y_pos += 0.35
        for g in short_goals:
            tb = slide.shapes.add_textbox(Inches(5.4), Inches(y_pos), Inches(4.1), Inches(0.3))
            tb.text_frame.text = f"• {truncate_text(g, 40)}"
            tb.text_frame.paragraphs[0].font.size = Pt(adjusted_font(DESIGN["fonts"]["body_small"], font_adj))
            tb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["text"])
            y_pos += 0.3
    y_pos += 0.15
    if medium_goals:
        tb = slide.shapes.add_textbox(Inches(5.2), Inches(y_pos), Inches(4.3), Inches(0.25))
        tb.text_frame.text = "Medium-Term (1-3 years)"
        tb.text_frame.paragraphs[0].font.size = Pt(adjusted_font(DESIGN["fonts"]["body_large"], font_adj))
        tb.text_frame.paragraphs[0].font.bold = True
        tb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["primary"])
        y_pos += 0.35
        for g in medium_goals:
            tb = slide.shapes.add_textbox(Inches(5.4), Inches(y_pos), Inches(4.1), Inches(0.3))
            tb.text_frame.text = f"• {truncate_text(g, 40)}"
            tb.text_frame.paragraphs[0].font.size = Pt(adjusted_font(DESIGN["fonts"]["body_small"], font_adj))
            tb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["text"])
            y_pos += 0.3

def render_market_position(slide, colors, data, page_num, layout_rec, context):
    font_adj = layout_rec.get("font_adjustment", 0)
    vertical = data.get("primaryVertical") or "technology"
    industry_content = get_industry_specific_content(vertical, "market-position")
    add_slide_header(slide, colors, "Market Position & Competitive Landscape", font_adj=font_adj)
    add_slide_footer(slide, colors, page_num)
    add_section_box(slide, colors, 0.3, 0.95, 4.5, 3.8, "Market Overview", font_adj=font_adj)
    tam = data.get("marketSize") or data.get("market_size") or "N/A"
    growth = data.get("marketGrowthRate") or data.get("market_growth_rate") or "N/A"
    add_metric_card(slide, colors, 0.45, 1.5, 4.2, 0.65, tam, "Total Addressable Market", font_adj)
    add_metric_card(slide, colors, 0.45, 2.35, 4.2, 0.65, f"{growth}%", "Market Growth Rate", font_adj)
    if industry_content.get("benchmarks_text"):
        tb = slide.shapes.add_textbox(Inches(0.45), Inches(3.2), Inches(4.2), Inches(0.4))
        tb.text_frame.text = industry_content["benchmarks_text"]
        tb.text_frame.paragraphs[0].font.size = Pt(adjusted_font(DESIGN["fonts"]["body_small"], font_adj))
        tb.text_frame.paragraphs[0].font.italic = True
        tb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["text_light"])
    add_section_box(slide, colors, 5.0, 0.95, 4.7, 3.8, "Competitive Advantages", colors["secondary"], font_adj)
    advantages = parse_pipe_separated(data.get("competitiveAdvantages") or data.get("competitive_advantages") or "", 5)
    y_pos = 1.5
    for adv in advantages:
        if adv:
            tb = slide.shapes.add_textbox(Inches(5.2), Inches(y_pos), Inches(4.3), Inches(0.4))
            tb.text_frame.text = f"• {truncate_text(adv[0] if len(adv)>0 else '', 35)}"
            tb.text_frame.paragraphs[0].font.size = Pt(adjusted_font(DESIGN["fonts"]["body"], font_adj))
            tb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["text"])
            y_pos += 0.5

def render_synergies(slide, colors, data, page_num, layout_rec, context):
    font_adj = layout_rec.get("font_adjustment", 0)
    buyer_types = data.get("targetBuyerType") or ["strategic"]
    add_slide_header(slide, colors, "Strategic Value & Synergies", font_adj=font_adj)
    add_slide_footer(slide, colors, page_num)
    if "strategic" in buyer_types:
        add_section_box(slide, colors, 0.3, 0.95, 4.5, 3.8, "Strategic Synergies", font_adj=font_adj)
        syns = parse_lines(data.get("synergiesStrategic") or data.get("synergies_strategic") or "", 6)
        y_pos = 1.5
        for s in syns:
            tb = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(4.1), Inches(0.4))
            tb.text_frame.text = f"• {truncate_text(s, 50)}"
            tb.text_frame.paragraphs[0].font.size = Pt(adjusted_font(DESIGN["fonts"]["body"], font_adj))
            tb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["text"])
            y_pos += 0.5
    if "financial" in buyer_types:
        add_section_box(slide, colors, 5.0, 0.95, 4.7, 3.8, "Financial Synergies", colors["secondary"], font_adj)
        fins = parse_lines(data.get("synergiesFinancial") or data.get("synergies_financial") or "", 6)
        y_pos = 1.5
        for s in fins:
            tb = slide.shapes.add_textbox(Inches(5.2), Inches(y_pos), Inches(4.3), Inches(0.4))
            tb.text_frame.text = f"• {truncate_text(s, 50)}"
            tb.text_frame.paragraphs[0].font.size = Pt(adjusted_font(DESIGN["fonts"]["body_small"], font_adj))
            tb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["text"])
            y_pos += 0.5

# ============================================================================
# TITLE & SPECIAL SLIDES
# ============================================================================
def render_title_slide(slide, colors, data, doc_config):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(5.625))
    bg.fill.solid(); bg.fill.fore_color.rgb = hex_to_rgb(colors["primary"]); bg.line.fill.background()
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.5), Inches(10), Inches(0.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = hex_to_rgb(colors["accent"]); bar.line.fill.background()
    company = data.get("companyName") or data.get("company_name") or "Company Name"
    codename = data.get("projectCodename") or data.get("project_codename") or "Project"
    doc_name = doc_config.get("name", "Information Memorandum")
    ttb = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(8), Inches(0.8))
    ttb.text_frame.text = company
    ttb.text_frame.paragraphs[0].font.size = Pt(48); ttb.text_frame.paragraphs[0].font.bold = True
    ttb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["white"])
    ttb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    stb = slide.shapes.add_textbox(Inches(1), Inches(2.7), Inches(8), Inches(0.5))
    stb.text_frame.text = doc_name
    stb.text_frame.paragraphs[0].font.size = Pt(24); stb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["white"])
    stb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    ctb = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.4))
    ctb.text_frame.text = f"Project {codename}"
    ctb.text_frame.paragraphs[0].font.size = Pt(18); ctb.text_frame.paragraphs[0].font.italic = True
    ctb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["white"])
    ctb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    dtb = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(8), Inches(0.3))
    dtb.text_frame.text = format_date(data.get("presentationDate") or data.get("presentation_date"))
    dtb.text_frame.paragraphs[0].font.size = Pt(14); dtb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["white"])
    dtb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    advisor = data.get("advisorName") or data.get("advisor_name") or ""
    if advisor:
        atb = slide.shapes.add_textbox(Inches(1), Inches(4.4), Inches(8), Inches(0.3))
        atb.text_frame.text = f"Prepared by {advisor}"
        atb.text_frame.paragraphs[0].font.size = Pt(12); atb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["white"])
        atb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def render_disclaimer_slide(slide, colors, data, page_num):
    add_slide_header(slide, colors, "Disclaimer"); add_slide_footer(slide, colors, page_num)
    add_section_box(slide, colors, 0.3, 0.95, 9.4, 3.8)
    txt = """This presentation has been prepared solely for informational purposes. The information contained herein is confidential and proprietary. By accepting this document, you agree to maintain its confidentiality and not to reproduce, distribute, or disclose it without prior written consent.\n\nThis presentation does not constitute an offer to sell or a solicitation to buy securities. Any investment decision should be made only after thorough due diligence and consultation with professional advisors.\n\nThe financial projections and forward-looking statements contained herein are based on assumptions that may or may not prove accurate. Actual results may vary materially."""
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9.0), Inches(3.2))
    tb.text_frame.text = txt
    tb.text_frame.paragraphs[0].font.size = Pt(11); tb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["text"])

def render_thank_you_slide(slide, colors, data, doc_config):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(5.625))
    bg.fill.solid(); bg.fill.fore_color.rgb = hex_to_rgb(colors["primary"]); bg.line.fill.background()
    ttb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(0.8))
    ttb.text_frame.text = "Thank You"
    ttb.text_frame.paragraphs[0].font.size = Pt(48); ttb.text_frame.paragraphs[0].font.bold = True
    ttb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["white"])
    ttb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    company = data.get("companyName") or data.get("company_name") or ""
    if company:
        ctb = slide.shapes.add_textbox(Inches(1), Inches(2.9), Inches(8), Inches(0.4))
        ctb.text_frame.text = company
        ctb.text_frame.paragraphs[0].font.size = Pt(20); ctb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["white"])
        ctb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    parts = []
    adv = data.get("advisorName") or data.get("advisor_name") or ""
    if adv: parts.append(f"Prepared by {adv}")
    em = data.get("contactEmail") or data.get("contact_email") or ""
    if em: parts.append(em)
    ph = data.get("contactPhone") or data.get("contact_phone") or ""
    if ph: parts.append(ph)
    if parts:
        ctb2 = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.8))
        ctb2.text_frame.text = "\n".join(parts)
        ctb2.text_frame.paragraphs[0].font.size = Pt(14); ctb2.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["white"])
        ctb2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    ftb = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(8), Inches(0.3))
    ftb.text_frame.text = "Strictly Private & Confidential"
    ftb.text_frame.paragraphs[0].font.size = Pt(10); ftb.text_frame.paragraphs[0].font.italic = True
    ftb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["white"])
    ftb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ============================================================================
# v8.3.0: NEW RENDERERS (TOC, Company Overview, Leadership, Risks)
# ============================================================================
def render_toc(slide, colors, data, page_num, layout_rec, context):
    font_adj = layout_rec.get("font_adjustment", 0)
    add_slide_header(slide, colors, "Table of Contents", font_adj=font_adj)
    add_slide_footer(slide, colors, page_num)
    labels = {"executive-summary":"Executive Summary","investment-highlights":"Investment Highlights",
        "company-overview":"Company Overview","services":"Service Lines & Capabilities",
        "clients":"Client Portfolio","financials":"Financial Performance",
        "case-study":"Case Studies","growth":"Growth Strategy & Roadmap",
        "market-position":"Market Position","synergies":"Strategic Value & Synergies",
        "risks":"Risk Factors & Mitigation","leadership":"Leadership Team"}
    doc_config = context.get("doc_config", {})
    entries = []
    for st in doc_config.get("required_slides", []) + doc_config.get("optional_slides", []):
        if st in labels: entries.append(labels[st])
    if not entries: entries = list(labels.values())
    col1 = entries[:len(entries)//2+1]; col2 = entries[len(entries)//2+1:]
    for ci, ents in enumerate([col1, col2]):
        cx = 0.5 + (ci*4.8); yp = 1.1
        for idx, e in enumerate(ents):
            sn = idx+1+(ci*len(col1))
            b = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(yp+0.02), Inches(0.35), Inches(0.35))
            b.fill.solid(); b.fill.fore_color.rgb = hex_to_rgb(colors["secondary"]); b.line.fill.background()
            ntb = slide.shapes.add_textbox(Inches(cx), Inches(yp+0.02), Inches(0.35), Inches(0.35))
            ntb.text_frame.paragraphs[0].text = f"{sn:02d}"
            ntb.text_frame.paragraphs[0].font.size = Pt(10); ntb.text_frame.paragraphs[0].font.bold = True
            ntb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["white"])
            ntb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            etb = slide.shapes.add_textbox(Inches(cx+0.5), Inches(yp), Inches(3.8), Inches(0.4))
            etb.text_frame.text = e
            etb.text_frame.paragraphs[0].font.size = Pt(adjusted_font(13, font_adj))
            etb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["text"])
            yp += 0.52

def render_company_overview(slide, colors, data, page_num, layout_rec, context):
    font_adj = layout_rec.get("font_adjustment", 0)
    add_slide_header(slide, colors, "Company Overview", font_adj=font_adj)
    add_slide_footer(slide, colors, page_num)
    add_section_box(slide, colors, 0.3, 0.95, 5.8, 2.4, "About the Company", font_adj=font_adj)
    desc = data.get("companyDescription") or data.get("company_description") or ""
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(1.45), Inches(5.5), Inches(1.7))
    tb.text_frame.text = truncate_description(desc, 400)
    tb.text_frame.paragraphs[0].font.size = Pt(adjusted_font(DESIGN["fonts"]["body"], font_adj))
    tb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["text"])
    tb.text_frame.word_wrap = True
    add_section_box(slide, colors, 6.3, 0.95, 3.4, 2.4, "Key Facts", colors["secondary"], font_adj)
    facts = []
    if data.get("foundedYear"): facts.append(("Founded", str(data["foundedYear"])))
    if data.get("headquarters"): facts.append(("Headquarters", truncate_text(str(data["headquarters"]),20)))
    if data.get("employeeCountFT"): facts.append(("Employees", str(data["employeeCountFT"])))
    if data.get("revenueFY25"): facts.append(("Revenue FY25", f"INR {data['revenueFY25']} Cr"))
    fy = 1.5
    for lbl, val in facts[:5]:
        lt = slide.shapes.add_textbox(Inches(6.45), Inches(fy), Inches(3.1), Inches(0.18))
        lt.text_frame.paragraphs[0].text = lbl
        lt.text_frame.paragraphs[0].font.size = Pt(9); lt.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["text_light"])
        vt = slide.shapes.add_textbox(Inches(6.45), Inches(fy+0.16), Inches(3.1), Inches(0.22))
        vt.text_frame.paragraphs[0].text = val
        vt.text_frame.paragraphs[0].font.size = Pt(12); vt.text_frame.paragraphs[0].font.bold = True
        vt.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["text"])
        fy += 0.42
    metrics = []
    if data.get("topClientCount") or data.get("totalClients"):
        metrics.append((str(data.get("topClientCount") or data.get("totalClients","")), "Total Clients", "●"))
    if data.get("top10Concentration"):
        metrics.append((f"{data['top10Concentration']}%", "Top 10 Conc.", "◆"))
    if data.get("netRetention"):
        metrics.append((f"{data['netRetention']}%", "Net Retention", "★"))
    if data.get("ebitdaMarginFY25"):
        metrics.append((f"{data['ebitdaMarginFY25']}%", "EBITDA Margin", "▲"))
    if metrics: add_metric_row(slide, colors, metrics, y=3.6, font_adj=font_adj)

def render_leadership(slide, colors, data, page_num, layout_rec, context):
    font_adj = layout_rec.get("font_adjustment", 0)
    add_slide_header(slide, colors, "Leadership Team", font_adj=font_adj)
    add_slide_footer(slide, colors, page_num)
    leadership_text = data.get("leadershipTeam") or data.get("leadership_team") or ""
    team = parse_pipe_separated(leadership_text, 8)
    if not team:
        fn = data.get("founderName") or ""
        if fn: team = [[fn, "Founder & CEO", data.get("founderBackground","")]]
    if not team: return
    # Tier 1: CEO card
    m = team[0]; name = m[0] if len(m)>0 else ""; title = m[1] if len(m)>1 else ""
    cw = 3.5; cx = (10-cw)/2
    psz = 0.7
    p = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx+(cw-psz)/2), Inches(1.05), Inches(psz), Inches(psz))
    p.fill.solid(); p.fill.fore_color.rgb = hex_to_rgb(colors["border"])
    p.line.color.rgb = hex_to_rgb(colors["primary"])
    inits = "".join([w[0].upper() for w in name.split()[:2]]) if name else ""
    itb = slide.shapes.add_textbox(Inches(cx+(cw-psz)/2), Inches(1.05), Inches(psz), Inches(psz))
    itb.text_frame.paragraphs[0].text = inits
    itb.text_frame.paragraphs[0].font.size = Pt(18); itb.text_frame.paragraphs[0].font.bold = True
    itb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["primary"])
    itb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    nc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx), Inches(1.85), Inches(cw), Inches(0.65))
    nc.fill.solid(); nc.fill.fore_color.rgb = hex_to_rgb(colors["primary"]); nc.line.fill.background()
    ntb = slide.shapes.add_textbox(Inches(cx+0.1), Inches(1.87), Inches(cw-0.2), Inches(0.3))
    ntb.text_frame.paragraphs[0].text = truncate_text(name, 35)
    ntb.text_frame.paragraphs[0].font.size = Pt(13); ntb.text_frame.paragraphs[0].font.bold = True
    ntb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["white"])
    ntb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    ttb = slide.shapes.add_textbox(Inches(cx+0.1), Inches(2.17), Inches(cw-0.2), Inches(0.25))
    ttb.text_frame.paragraphs[0].text = truncate_text(title, 35)
    ttb.text_frame.paragraphs[0].font.size = Pt(10); ttb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["white"])
    ttb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Tier 2
    rem = team[1:7]
    if rem:
        num = len(rem); cols = min(num,3); gap = 0.2; cw2 = 2.8
        tw = (cols*cw2)+((cols-1)*gap); sx = (10-tw)/2; sy = 2.75
        for idx, m in enumerate(rem):
            r = idx//cols; c = idx%cols
            mx = sx + c*(cw2+gap); my = sy + r*1.2
            n2 = m[0] if len(m)>0 else ""; t2 = m[1] if len(m)>1 else ""
            cd = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(mx), Inches(my), Inches(cw2), Inches(0.95))
            cd.fill.solid(); cd.fill.fore_color.rgb = hex_to_rgb(colors["light_bg"])
            cd.line.color.rgb = hex_to_rgb(colors["border"])
            ac = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(mx), Inches(my), Inches(0.06), Inches(0.95))
            ac.fill.solid(); ac.fill.fore_color.rgb = hex_to_rgb(colors["secondary"]); ac.line.fill.background()
            nt = slide.shapes.add_textbox(Inches(mx+0.15), Inches(my+0.1), Inches(cw2-0.25), Inches(0.3))
            nt.text_frame.paragraphs[0].text = truncate_text(n2,30)
            nt.text_frame.paragraphs[0].font.size = Pt(11); nt.text_frame.paragraphs[0].font.bold = True
            nt.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["text"])
            tt = slide.shapes.add_textbox(Inches(mx+0.15), Inches(my+0.4), Inches(cw2-0.25), Inches(0.45))
            tt.text_frame.paragraphs[0].text = truncate_text(t2,40)
            tt.text_frame.paragraphs[0].font.size = Pt(9); tt.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["text_light"])
            tt.text_frame.word_wrap = True

def render_risk_factors(slide, colors, data, page_num, layout_rec, context):
    font_adj = layout_rec.get("font_adjustment", 0)
    add_slide_header(slide, colors, "Risk Factors & Mitigation", font_adj=font_adj)
    add_slide_footer(slide, colors, page_num)
    risks = []
    if data.get("businessRisks"): risks.append(("Business Risks", data["businessRisks"]))
    if data.get("marketRisks"): risks.append(("Market Risks", data["marketRisks"]))
    if data.get("operationalRisks"): risks.append(("Operational Risks", data["operationalRisks"]))
    if data.get("mitigationStrategies"): risks.append(("Mitigation Strategies", data["mitigationStrategies"]))
    if not risks and data.get("riskFactors"):
        items = parse_lines(data["riskFactors"], 8)
        if items: risks.append(("Key Risk Factors", "\n".join(items)))
    if not risks: return
    num = len(risks)
    if num <= 2:
        for i, (cat, content) in enumerate(risks):
            bx = 0.3 + (i*4.85); bw = 4.55
            hc = colors["accent"] if "Mitigation" in cat else colors["primary"]
            add_section_box(slide, colors, bx, 0.95, bw, 3.8, cat, hc, font_adj)
            tb = slide.shapes.add_textbox(Inches(bx+0.15), Inches(1.45), Inches(bw-0.3), Inches(3.1))
            tb.text_frame.text = truncate_description(content, 350)
            tb.text_frame.paragraphs[0].font.size = Pt(adjusted_font(DESIGN["fonts"]["body_small"], font_adj))
            tb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["text"])
            tb.text_frame.word_wrap = True
    else:
        rh = min(0.9, 3.8/num - 0.05); yp = 0.95
        hcs = [colors["primary"], colors["accent"], colors["secondary"], colors.get("warning","D69E2E")]
        for i, (cat, content) in enumerate(risks):
            hc = hcs[i % len(hcs)]
            hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(yp), Inches(9.4), Inches(0.3))
            hdr.fill.solid(); hdr.fill.fore_color.rgb = hex_to_rgb(hc); hdr.line.fill.background()
            htb = slide.shapes.add_textbox(Inches(0.42), Inches(yp+0.02), Inches(9.1), Inches(0.26))
            htb.text_frame.paragraphs[0].text = cat
            htb.text_frame.paragraphs[0].font.size = Pt(adjusted_font(11, font_adj)); htb.text_frame.paragraphs[0].font.bold = True
            htb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["white"])
            cbg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(yp+0.3), Inches(9.4), Inches(rh-0.3))
            cbg.fill.solid(); cbg.fill.fore_color.rgb = hex_to_rgb(colors["light_bg"])
            cbg.line.color.rgb = hex_to_rgb(colors["border"])
            ctb = slide.shapes.add_textbox(Inches(0.45), Inches(yp+0.35), Inches(9.1), Inches(rh-0.4))
            ctb.text_frame.text = truncate_description(content, 180)
            ctb.text_frame.paragraphs[0].font.size = Pt(adjusted_font(10, font_adj))
            ctb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["text"])
            ctb.text_frame.word_wrap = True
            yp += rh + 0.05

# ============================================================================
# APPENDIX SLIDES
# ============================================================================
def render_appendix_financials(slide, colors, data, page_num, layout_rec, context):
    font_adj = layout_rec.get("font_adjustment", -1)
    add_slide_header(slide, colors, "Appendix A: Detailed Financial Information", font_adj=font_adj)
    add_slide_footer(slide, colors, page_num)
    add_section_box(slide, colors, 0.3, 0.95, 9.4, 3.8, "Financial Details")
    yp = 1.5
    for st, content in [("Revenue Breakdown by Service Line", data.get("revenueByService") or ""),
                        ("Cost Structure Analysis", data.get("costStructure") or ""),
                        ("Working Capital Requirements", data.get("workingCapital") or "")]:
        if content:
            tt = slide.shapes.add_textbox(Inches(0.5), Inches(yp), Inches(9.0), Inches(0.25))
            tt.text_frame.text = st
            tt.text_frame.paragraphs[0].font.size = Pt(adjusted_font(12, font_adj))
            tt.text_frame.paragraphs[0].font.bold = True
            tt.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["primary"])
            ct = slide.shapes.add_textbox(Inches(0.5), Inches(yp+0.3), Inches(9.0), Inches(0.6))
            ct.text_frame.text = truncate_description(content, 300)
            ct.text_frame.paragraphs[0].font.size = Pt(adjusted_font(10, font_adj))
            ct.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["text"])
            yp += 1.0

def render_appendix_case_studies(slide, colors, data, page_num, layout_rec, context):
    font_adj = layout_rec.get("font_adjustment", -1)
    add_slide_header(slide, colors, "Appendix B: Additional Case Studies", font_adj=font_adj)
    add_slide_footer(slide, colors, page_num)
    add_section_box(slide, colors, 0.3, 0.95, 9.4, 3.8)
    cs = data.get("caseStudies") or []; extra = cs[2:] if len(cs)>2 else []
    yp = 1.3
    for s in extra[:2]:
        cl = s.get("client","Client")
        tt = slide.shapes.add_textbox(Inches(0.5), Inches(yp), Inches(9.0), Inches(0.25))
        tt.text_frame.text = f"Case Study: {truncate_text(cl, 60)}"
        tt.text_frame.paragraphs[0].font.size = Pt(adjusted_font(12, font_adj)); tt.text_frame.paragraphs[0].font.bold = True
        tt.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["primary"])
        yp += 0.3
        c = f"Challenge: {truncate_text(s.get('challenge',''),100)} | Solution: {truncate_text(s.get('solution',''),100)} | Results: {truncate_text(s.get('results',''),100)}"
        ct = slide.shapes.add_textbox(Inches(0.5), Inches(yp), Inches(9.0), Inches(0.6))
        ct.text_frame.text = c
        ct.text_frame.paragraphs[0].font.size = Pt(adjusted_font(9, font_adj))
        ct.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["text"])
        yp += 0.9

def render_appendix_team_bios(slide, colors, data, page_num, layout_rec, context):
    font_adj = layout_rec.get("font_adjustment", -1)
    add_slide_header(slide, colors, "Appendix C: Detailed Team Biographies", font_adj=font_adj)
    add_slide_footer(slide, colors, page_num)
    add_section_box(slide, colors, 0.3, 0.95, 9.4, 3.8)
    lt = data.get("leadershipTeam") or data.get("leadership_team") or ""
    team = parse_pipe_separated(lt, 4)
    yp = 1.3
    for m in team:
        if m and len(m) >= 2:
            tb = slide.shapes.add_textbox(Inches(0.5), Inches(yp), Inches(9.0), Inches(0.4))
            tb.text_frame.text = f"• {m[0]} — {m[1]}"
            tb.text_frame.paragraphs[0].font.size = Pt(adjusted_font(11, font_adj))
            tb.text_frame.paragraphs[0].font.bold = True
            tb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["text"])
            yp += 0.5

# ============================================================================
# UNIVERSAL createSlide() WRAPPER
# ============================================================================
def create_slide(slide_type: str, prs: Presentation, colors: dict, data: dict, page_num: int, context: dict) -> Optional[int]:
    # Guard clauses
    if slide_type == "case-study" and not data.get("caseStudies") and not data.get("cs1Client"): return None
    if slide_type == "financials" and not (data.get("revenueFY24") or data.get("revenueFY25")): return None
    if slide_type == "market-position" and not (data.get("marketSize") or data.get("competitiveAdvantages")): return None
    if slide_type == "synergies" and not (data.get("synergiesStrategic") or data.get("synergiesFinancial")): return None
    if slide_type == "appendix-team-bios" and not data.get("teamBios") and not data.get("leadershipTeam"): return None
    if slide_type == "appendix-case-studies":
        if len(data.get("caseStudies") or []) <= 2: return None

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    layout_rec = analyze_data_for_layout_sync(data, slide_type)

    if slide_type == "title":
        render_title_slide(slide, colors, data, context.get("doc_config", {})); return None
    elif slide_type == "disclaimer":
        render_disclaimer_slide(slide, colors, data, page_num); return page_num + 1
    elif slide_type == "toc":
        render_toc(slide, colors, data, page_num, layout_rec, context); return page_num + 1
    elif slide_type == "executive-summary":
        render_executive_summary(slide, colors, data, page_num, layout_rec, context); return page_num + 1
    elif slide_type == "investment-highlights":
        add_slide_header(slide, colors, "Investment Highlights"); add_slide_footer(slide, colors, page_num)
        hl = parse_lines(data.get("investmentHighlights") or "", 8)
        yp = 1.1
        for i, h in enumerate(hl):
            c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(yp+0.02), Inches(0.28), Inches(0.28))
            c.fill.solid(); c.fill.fore_color.rgb = hex_to_rgb(colors["secondary"]); c.line.fill.background()
            itb = slide.shapes.add_textbox(Inches(0.5), Inches(yp+0.02), Inches(0.28), Inches(0.28))
            itb.text_frame.paragraphs[0].text = "✦"
            itb.text_frame.paragraphs[0].font.size = Pt(10)
            itb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["white"])
            itb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            tb = slide.shapes.add_textbox(Inches(0.9), Inches(yp), Inches(8.8), Inches(0.35))
            tb.text_frame.text = truncate_text(h, 85)
            tb.text_frame.paragraphs[0].font.size = Pt(12)
            tb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(colors["text"])
            yp += 0.45
        return page_num + 1
    elif slide_type == "company-overview":
        render_company_overview(slide, colors, data, page_num, layout_rec, context); return page_num + 1
    elif slide_type == "services":
        render_services(slide, colors, data, page_num, layout_rec, context); return page_num + 1
    elif slide_type == "clients":
        render_clients(slide, colors, data, page_num, layout_rec, context); return page_num + 1
    elif slide_type == "financials":
        render_financials(slide, colors, data, page_num, layout_rec, context); return page_num + 1
    elif slide_type == "case-study":
        cs = data.get("caseStudies") or []
        if not cs and data.get("cs1Client"):
            cs.append({"client": data.get("cs1Client"), "challenge": data.get("cs1Challenge"),
                        "solution": data.get("cs1Solution"), "results": data.get("cs1Results")})
        if cs: render_case_study(slide, colors, data, page_num, cs[0], layout_rec, context); return page_num + 1
        return None
    elif slide_type == "growth":
        render_growth(slide, colors, data, page_num, layout_rec, context); return page_num + 1
    elif slide_type == "market-position":
        render_market_position(slide, colors, data, page_num, layout_rec, context); return page_num + 1
    elif slide_type == "leadership":
        render_leadership(slide, colors, data, page_num, layout_rec, context); return page_num + 1
    elif slide_type == "synergies":
        render_synergies(slide, colors, data, page_num, layout_rec, context); return page_num + 1
    elif slide_type == "risks":
        render_risk_factors(slide, colors, data, page_num, layout_rec, context); return page_num + 1
    elif slide_type == "section-divider":
        st = data.get("_section_divider_title", "Section")
        sn = data.get("_section_divider_num")
        render_section_divider(slide, colors, st, sn); return None
    elif slide_type == "appendix-financials":
        render_appendix_financials(slide, colors, data, page_num, layout_rec, context); return page_num + 1
    elif slide_type == "appendix-case-studies":
        render_appendix_case_studies(slide, colors, data, page_num, layout_rec, context); return page_num + 1
    elif slide_type == "appendix-team-bios":
        render_appendix_team_bios(slide, colors, data, page_num, layout_rec, context); return page_num + 1
    elif slide_type == "thank-you":
        render_thank_you_slide(slide, colors, data, context.get("doc_config", {})); return None
    else:
        return None

# ============================================================================
# MAIN GENERATOR
# ============================================================================
def generate_presentation(data: Dict, theme: str = "modern-blue") -> Presentation:
    if isinstance(data, str):
        import json
        try: data = json.loads(data)
        except: data = {}
    if not isinstance(data, dict): data = {}
    if not data.get("companyName") and not data.get("company_name"): data["companyName"] = "Company Name"
    if not data.get("documentType") and not data.get("document_type"): data["documentType"] = "management-presentation"
    try:
        prs = Presentation(); prs.slide_width = Inches(10); prs.slide_height = Inches(5.625)
    except Exception as e:
        print(f"ERROR: Failed to create presentation: {e}"); raise
    try: colors = get_theme_colors(theme)
    except: colors = get_theme_colors("modern-blue")
    doc_type = (data.get("documentType") or data.get("document_type") or "management-presentation").lower()
    if doc_type not in ["management-presentation","cim","teaser"]: doc_type = "management-presentation"
    doc_config = DOCUMENT_CONFIGS.get(doc_type, DOCUMENT_CONFIGS["management-presentation"])
    primary_vertical = (data.get("primaryVertical") or data.get("primary_vertical") or "technology").lower()
    industry_data = INDUSTRY_DATA.get(primary_vertical, INDUSTRY_DATA.get("technology", {}))
    context = {"doc_config": doc_config, "industry_data": industry_data,
               "buyer_types": data.get("targetBuyerType") or data.get("target_buyer_type") or ["strategic"]}
    try: slides_to_generate = get_slides_for_document_type(doc_type, data)
    except Exception as e:
        print(f"ERROR: Failed to determine slides: {e}")
        import traceback; traceback.print_exc()
        slides_to_generate = ["title","disclaimer","executive-summary","services","clients","financials","thank-you"]
    print(f"=== GENERATION SUMMARY (v8.3.0) ===")
    print(f"Document Type: {doc_type} | Industry: {primary_vertical} | Theme: {theme}")
    print(f"Slides ({len(slides_to_generate)}): {slides_to_generate}")
    print(f"====================================")
    page_num = 1; slides_created = 0
    for slide_type in slides_to_generate:
        try:
            result = create_slide(slide_type, prs, colors, data, page_num, context)
            if result is not None: page_num = result
            slides_created += 1
            print(f"✓ Created slide: {slide_type}")
        except Exception as e:
            print(f"✗ ERROR creating slide '{slide_type}': {e}")
            import traceback; traceback.print_exc()
            continue
    print(f"=== GENERATION COMPLETE (v8.3.0) ===")
    print(f"Total slides created: {slides_created}/{len(slides_to_generate)}")
    print(f"====================================")
    return prs
